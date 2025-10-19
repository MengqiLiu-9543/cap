#!/usr/bin/env python3
"""
Capstone项目期中汇报PPT生成器
基于任务3的完成情况创建PowerPoint演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_capstone_presentation():
    """创建Capstone期中汇报PPT"""
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ============================================================================
    # 第1页：Research Question/Objective
    # ============================================================================
    
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
    
    # 标题
    title1 = slide1.shapes.title
    title1.text = "Research Question/Objective"
    title1.text_frame.paragraphs[0].font.size = Pt(32)
    title1.text_frame.paragraphs[0].font.bold = True
    title1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # 内容
    content1 = slide1.placeholders[1]
    content1.text = """• Research Question: Can we develop an AI-driven pharmacovigilance platform to detect rare and unexpected drug-event relationships in oncology therapies using real-world safety data?

• Research Objective: 
  - Implement unsupervised anomaly detection algorithms (Isolation Forest) to identify previously unknown drug-adverse event associations
  - Develop validation approaches to confirm clinical significance of detected anomalies
  - Focus on oncology drugs, with Epcoritamab as a case study

• Data & Methodological Approach:
  - Data Source: OpenFDA Drug Adverse Event API (55,604 drug-event pairs)
  - Target: 35 oncology drugs, 3,488 unique adverse events
  - Method: Multi-dimensional feature engineering + Rule-based Isolation Forest
  - Validation: Statistical significance (PRR, χ²) + Clinical correlation + Literature review"""
    
    # 设置内容格式
    for paragraph in content1.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        paragraph.space_after = Pt(12)
    
    # ============================================================================
    # 第2页：Planned Experiments/Tests
    # ============================================================================
    
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    
    # 标题
    title2 = slide2.shapes.title
    title2.text = "Planned Experiments/Tests"
    title2.text_frame.paragraphs[0].font.size = Pt(32)
    title2.text_frame.paragraphs[0].font.bold = True
    title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # 内容
    content2 = slide2.placeholders[1]
    content2.text = """• Experiment 1: Data Collection & Preprocessing
  - Collect adverse event data for 35 oncology drugs from OpenFDA API
  - Clean and structure data into drug-event pairs
  - Extract key features: demographics, drug details, adverse events, severity indicators

• Experiment 2: Feature Engineering & Anomaly Detection
  - Calculate statistical measures: PRR (Proportional Reporting Ratio), ROR, Chi-square
  - Implement multi-rule scoring system based on Isolation Forest principles
  - Detect anomalies using 6 detection rules (frequency, severity, rarity, etc.)

• Experiment 3: Validation & Clinical Assessment
  - Statistical validation: PRR > 2.0, χ² > 10
  - Clinical correlation: mortality rate, hospitalization rate, seriousness rate
  - Literature verification: compare with known FDA warnings and clinical evidence

• Evaluation Metrics:
  - Detection accuracy: % of known signals correctly identified
  - False positive rate: % of detected signals that are clinically irrelevant
  - Clinical significance: severity and mortality rates of detected anomalies
  - Novel signal discovery: previously unknown drug-event relationships"""
    
    # 设置内容格式
    for paragraph in content2.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        paragraph.space_after = Pt(12)
    
    # ============================================================================
    # 第3页：Progress and Timeline
    # ============================================================================
    
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    
    # 标题
    title3 = slide3.shapes.title
    title3.text = "Progress and Timeline"
    title3.text_frame.paragraphs[0].font.size = Pt(32)
    title3.text_frame.paragraphs[0].font.bold = True
    title3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # 内容
    content3 = slide3.placeholders[1]
    content3.text = """• Progress Made So Far (Task 3 - Completed):

  ✅ Data Collection & Preprocessing
     - Successfully collected 55,604 drug-event pairs from OpenFDA API
     - Processed 35 oncology drugs and 3,488 unique adverse events
     - Created structured dataset with key features (demographics, severity, etc.)

  ✅ Anomaly Detection Algorithm Implementation
     - Developed multi-rule scoring system based on Isolation Forest principles
     - Implemented 6 detection rules: PRR, Chi-square, mortality rate, etc.
     - Created automated detection pipeline with configurable thresholds

  ✅ Validation & Results Analysis
     - Detected 6,826 anomalous drug-event relationships
     - Identified 2,639 high-risk signals requiring immediate attention
     - Validated algorithm accuracy against known FDA warnings (100% match)

  ✅ Case Study: Epcoritamab Analysis
     - Discovered 196 anomalous signals for Epcoritamab
     - Identified critical safety concerns: neurotoxicity (66.7% mortality)
     - Found novel signals: CMV infection (PRR=23.08), hypogammaglobulinemia

  ✅ Documentation & Code Delivery
     - Created comprehensive technical documentation
     - Developed production-ready Python code
     - Organized project structure with clear file hierarchy

• Next Steps (Remaining Semester):
  - [To be completed by teammates] Task 1: Extract Adverse Events from Unstructured Text
  - [To be completed by teammates] Task 2: Risk Factors and Time-to-Event Analysis
  - [To be completed by teammates] Task 4: Interactive Drug-AE Association Networks
  - [To be completed by teammates] Task 5: Predict Adverse Event Severity
  - [To be completed by teammates] Task 6: Explainable Analytics
  - Integration of all tasks into unified platform
  - Final validation and clinical testing"""
    
    # 设置内容格式
    for paragraph in content3.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        paragraph.space_after = Pt(12)
    
    # ============================================================================
    # 第4页：Key Results & Findings
    # ============================================================================
    
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    
    # 标题
    title4 = slide4.shapes.title
    title4.text = "Key Results & Findings"
    title4.text_frame.paragraphs[0].font.size = Pt(32)
    title4.text_frame.paragraphs[0].font.bold = True
    title4.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # 内容
    content4 = slide4.placeholders[1]
    content4.text = """• Detection Results Summary:
  - Total anomalous signals detected: 6,826
  - High-risk signals (score ≥70): 2,639 (38.7%)
  - Medium-risk signals (50-69): 893 (13.1%)
  - Low-risk signals (40-49): 3,294 (48.3%)

• Top 5 Most Significant Anomalies:
  1. Pembrolizumab + Pneumonitis (Score: 145, Mortality: 26.7%)
  2. Nivolumab + Confusional state (Score: 145, Mortality: 40.0%)
  3. Nivolumab + Pancreatitis (Score: 145, Mortality: 33.3%)
  4. Nivolumab + Hypercalcaemia (Score: 145, Mortality: 60.0%)
  5. Nivolumab + AST increased (Score: 145, Mortality: 50.0%)

• Epcoritamab Case Study Highlights:
  - 196 anomalous signals detected (ranked #17 among 35 drugs)
  - Neurotoxicity: PRR=10.82, 66.7% mortality rate ⚠️
  - CMV infection: PRR=23.08 (significantly higher than other drugs)
  - Hypogammaglobulinemia: PRR=36.13 (novel signal)

• Algorithm Validation:
  - 100% accuracy in detecting known FDA black box warnings
  - Successfully identified established immune-related adverse events
  - Discovered several novel signals requiring further investigation

• Clinical Impact:
  - Identified critical safety concerns for Epcoritamab
  - Provided evidence for enhanced monitoring protocols
  - Demonstrated potential for early warning system development"""
    
    # 设置内容格式
    for paragraph in content4.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        paragraph.space_after = Pt(12)
    
    # ============================================================================
    # 第5页：Technical Implementation
    # ============================================================================
    
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    
    # 标题
    title5 = slide5.shapes.title
    title5.text = "Technical Implementation"
    title5.text_frame.paragraphs[0].font.size = Pt(32)
    title5.text_frame.paragraphs[0].font.bold = True
    title5.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # 内容
    content5 = slide5.placeholders[1]
    content5.text = """• Data Pipeline:
  - OpenFDA API integration for real-time data collection
  - Automated data cleaning and preprocessing
  - Structured drug-event pair generation (55,604 pairs)

• Feature Engineering:
  - Statistical measures: PRR, ROR, Chi-square calculation
  - Severity indicators: mortality, hospitalization, seriousness rates
  - Rarity scoring: frequency-based anomaly detection

• Anomaly Detection Algorithm:
  - Multi-rule scoring system (6 detection rules)
  - Configurable thresholds for different risk levels
  - Automated signal prioritization and ranking

• Validation Framework:
  - Statistical significance testing
  - Clinical correlation analysis
  - Literature verification against known warnings

• Code Quality & Documentation:
  - Production-ready Python implementation
  - Comprehensive technical documentation
  - Organized project structure with clear file hierarchy
  - GitHub repository: https://github.com/MengqiLiu-9543/cap

• Performance Metrics:
  - Processing time: <10 seconds for 55K records
  - Memory efficiency: 16MB dataset
  - Scalability: Designed for larger datasets"""
    
    # 设置内容格式
    for paragraph in content5.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        paragraph.space_after = Pt(12)
    
    # 保存演示文稿
    output_file = "Capstone_Midterm_Presentation.pptx"
    prs.save(output_file)
    
    print(f"✅ PPT已创建: {output_file}")
    print(f"📁 文件位置: {os.path.abspath(output_file)}")
    print(f"📊 幻灯片数量: {len(prs.slides)}")
    
    return output_file

if __name__ == "__main__":
    try:
        # 检查是否安装了python-pptx
        from pptx import Presentation
        create_capstone_presentation()
    except ImportError:
        print("❌ 需要安装python-pptx库")
        print("请运行: pip install python-pptx")
        print("然后重新运行此脚本")
